"""用于早期集成的结构化沙箱资产生成。"""

from __future__ import annotations

import json
import base64
from dataclasses import asdict
from pathlib import Path
from typing import Any

from pydantic import BaseModel, ConfigDict, Field

from src.ai_modules.config import get_sandbox_root, get_settings
from src.ai_modules.generation.content_chain import (
    ContentGenerationChain,
    GeneratedCodeAsset,
    GeneratedMindMap,
    GeneratedSectionBundle,
    GeneratedSlideDeck,
    GeneratedTextAsset,
)
from src.ai_modules.models import (
    VideoGenerationTaskPayload,
    VideoSandboxArtifact,
)
from src.ai_modules.runtime import SystemSnapshot


class GeneratedAsset(BaseModel):
    """写入沙箱存储的已生成资产的元数据。"""

    asset_type: str = Field(alias="assetType")
    title: str
    summary: str
    display_mode: str = Field(alias="displayMode")
    file_name: str = Field(default="", alias="fileName")
    local_path: str | None = Field(default=None, alias="localPath")
    preview_text: str = Field(alias="previewText")
    mime_type: str | None = Field(default=None, alias="mimeType")
    inline_content: str | None = Field(default=None, alias="inlineContent")
    language: str | None = None
    explanation: str | None = None
    thumbnail_path: str | None = Field(default=None, alias="thumbnailPath")
    thumbnail_file_name: str | None = Field(default=None, alias="thumbnailFileName")
    thumbnail_mime_type: str | None = Field(default=None, alias="thumbnailMimeType")
    duration_seconds: int | None = Field(default=None, alias="durationSeconds")
    video_style: str | None = Field(default=None, alias="videoStyle")
    knowledge_point: str | None = Field(default=None, alias="knowledgePoint")
    generated_by: str | None = Field(default=None, alias="generatedBy")
    content_origin: str | None = Field(default=None, alias="contentOrigin")
    provider: str | None = None
    model: str | None = None
    agent_name: str | None = Field(default=None, alias="agentName")
    evidence_ids: list[str] = Field(default_factory=list, alias="evidenceIds")
    fallback: bool | None = None
    from_cache: bool = Field(default=False, alias="fromCache")

    model_config = ConfigDict(populate_by_name=True)


class SectionPlan(BaseModel):
    """已生成教学资产中的规划章节。"""

    title: str
    objective: str
    source_titles: list[str] = Field(default_factory=list, alias="sourceTitles")

    model_config = ConfigDict(populate_by_name=True)


class ResourceGenerationService:
    """将结构化资产输出写入本地沙箱目录。"""

    def __init__(
        self,
        sandbox_root: Path | None = None,
        content_chain: ContentGenerationChain | None = None,
    ) -> None:
        self.sandbox_root = sandbox_root or get_sandbox_root()
        self.content_chain = content_chain or ContentGenerationChain()

    def build_asset(
        self,
        *,
        asset_type: str,
        params: dict,
        snapshot: SystemSnapshot,
    ) -> GeneratedAsset:
        builder_map = {
            "DOCUMENT": self._build_document,
            "READING": self._build_reading,
            "SLIDES": self._build_slides,
            "MINDMAP": self._build_mindmap,
            "CODE": self._build_code,
            "VIDEO": self._build_video,
        }
        builder = builder_map.get(asset_type)
        if builder is None:
            raise ValueError(f"Unsupported assetType: {asset_type}")
        return builder(params=params, snapshot=snapshot)

    async def build_video_asset(
        self,
        *,
        params: dict,
        snapshot: SystemSnapshot,
    ) -> GeneratedAsset:
        topic = str(params.get("topic") or params.get("query") or "教学主题")
        task_id = str(params.get("taskId") or "video-task")
        style = str(params.get("style") or "hybrid")
        duration_target = self._normalize_duration_seconds(params.get("duration") or params.get("durationTarget"))
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", []) if isinstance(retrieval, dict) else []
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        script_payload = await self.content_chain.generate_video_script_async(
            title=f"{topic}教学视频",
            topic=topic,
            snapshot=generation_snapshot,
            sources=sources,
            duration_seconds=duration_target,
            style=style,
            fallback_builder=self,
        )
        task_dir = self.sandbox_root / f"video_{self._safe_task_id(task_id)}"
        task_dir.mkdir(parents=True, exist_ok=True)

        script_json_path = task_dir / "script.json"
        script_text_path = task_dir / "script.txt"
        audio_path = task_dir / "speech.mp3"
        thumbnail_path = task_dir / "thumbnail.svg"

        script_json_path.write_text(
            json.dumps(script_payload.model_dump(by_alias=True), ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        script_text_path.write_text(script_payload.full_text, encoding="utf-8")

        tts_audio_bytes = params.get("tts_audio_bytes")
        if not tts_audio_bytes or not isinstance(tts_audio_bytes, bytes) or len(tts_audio_bytes) < 100:
            from src.ai_modules.llms.mimo_client import MiMoClient

            try:
                mimo_client = MiMoClient()
                tts_audio_bytes = await mimo_client.synthesize_speech(
                    text=script_payload.full_text[:1600],
                    style_description="用清晰自然的语速播报，声音沉稳专业，适合教学场景",
                    voice="mimo_default",
                    audio_format="mp3",
                )
            except Exception as exc:
                raise RuntimeError(
                    f"Video TTS generation failed: {type(exc).__name__}: {exc}"
                ) from exc
            params["tts_audio_bytes"] = tts_audio_bytes
        audio_path.write_bytes(tts_audio_bytes)

        thumbnail_path.write_text(
            self._build_video_thumbnail_svg(
                title=script_payload.title,
                topic=topic,
                style=style,
            ),
            encoding="utf-8",
        )

        artifact = VideoSandboxArtifact(
            taskDir=str(task_dir),
            scriptJsonPath=str(script_json_path),
            scriptTextPath=str(script_text_path),
            audioPath=str(audio_path),
            finalVideoPath=None,
            thumbnailPath=str(thumbnail_path),
            durationSeconds=script_payload.total_duration,
            videoStyle=style,
            previewText=script_payload.full_text[:100],
            summaryText=f"{topic} 教学视频脚本与语音已生成，等待浏览器本地渲染。",
            audioBase64=base64.b64encode(tts_audio_bytes).decode("utf-8"),
            audioFormat="mp3",
            avatarDataUrl="/dh_live/assets/combined_data.json.gz",
        )
        settings = get_settings()
        params["videoSandboxArtifact"] = artifact.model_dump(by_alias=True)
        params["videoGenerationTask"] = VideoGenerationTaskPayload(
            status="completed",
            title=script_payload.title,
            topic=topic,
            script=script_payload,
            durationSeconds=artifact.duration_seconds,
            videoStyle=style,
            ttsProvider=settings.tts_provider,
            avatarProvider="browser_dh_live_mini",
            generationParams={
                "durationTarget": script_payload.total_duration,
                "style": style,
            },
        ).model_dump(by_alias=True)
        return GeneratedAsset(
            assetType="VIDEO",
            title=script_payload.title,
            summary=artifact.summary_text,
            displayMode="VIDEO_PLAYER",
            fileName="browser-rendered.webm",
            localPath=None,
            previewText=artifact.preview_text,
            mimeType="video/webm",
            thumbnailPath=str(thumbnail_path),
            thumbnailFileName="thumbnail.svg",
            thumbnailMimeType="image/svg+xml",
            durationSeconds=artifact.duration_seconds,
            videoStyle=style,
            knowledgePoint=topic,
        )

    def _build_document(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        title = f"{params.get('query', '学习资源')}导学文档"
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", [])
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        section_plans = self._plan_document_sections(
            params=params,
            snapshot=snapshot,
            sources=sources,
        )
        generated_sections = self.content_chain.generate_document_sections(
            title=title,
            topic=str(params.get("rewrittenQuery", params.get("query", "主题"))),
            snapshot=generation_snapshot,
            section_plans=[plan.model_dump(by_alias=True) for plan in section_plans],
            sources=sources,
            fallback_builder=self,
        )
        content = self._render_document_markdown(
            title=title,
            params=params,
            snapshot=snapshot,
            section_plans=section_plans,
            sources=sources,
            generated_sections=generated_sections,
        )
        file_name = self._scoped_file_name("document_guide", "md", params)
        path = self._write_text(file_name, content)
        return GeneratedAsset(
            assetType="DOCUMENT",
            title=title,
            summary="基于检索证据生成的结构化课程导学文档",
            displayMode="MARKDOWN_CARD",
            fileName=file_name,
            localPath=str(path),
            previewText=self._build_preview_text(section_plans),
            mimeType="text/markdown; charset=UTF-8",
            inlineContent=content,
        )

    def _plan_document_sections(
        self,
        *,
        params: dict,
        snapshot: SystemSnapshot,
        sources: list[dict[str, Any]],
    ) -> list[SectionPlan]:
        topic = str(params.get("rewrittenQuery", params.get("query", "主题")))
        source_titles = [str(item.get("title", "未知来源")) for item in sources[:5]]
        learner_gap = ", ".join(snapshot.knowledge_gaps[:2]) or "暂无明确薄弱点"
        return [
            SectionPlan(
                title="一、核心概念与学习目标",
                objective=f"帮助学生围绕 `{topic}` 建立基础概念框架。",
                sourceTitles=source_titles[:2],
            ),
            SectionPlan(
                title="二、关键原理与判断方法",
                objective=f"从 `{snapshot.current_course}` 视角解释原理与判断依据。",
                sourceTitles=source_titles[:3],
            ),
            SectionPlan(
                title="三、典型误区与辨析",
                objective=f"聚焦薄弱点 `{learner_gap}`，说明常见误区与纠偏方式。",
                sourceTitles=source_titles[1:4],
            ),
            SectionPlan(
                title="四、练习建议与复习路径",
                objective="给出可执行的练习顺序、复习策略和下一步建议。",
                sourceTitles=source_titles[:2],
            ),
        ]

    def _render_document_markdown(
        self,
        *,
        title: str,
        params: dict,
        snapshot: SystemSnapshot,
        section_plans: list[SectionPlan],
        sources: list[dict[str, Any]],
        generated_sections: GeneratedSectionBundle,
    ) -> str:
        topic = str(params.get("rewrittenQuery", params.get("query", "主题")))
        lines = [
            f"# {title}",
            "",
            f"- 课程: {snapshot.current_course}",
            f"- 章节: {snapshot.current_chapter}",
            f"- 学生水平: {snapshot.student_level}",
            f"- 学习风格: {snapshot.preferred_style}",
            "",
            "## 文档概览",
            f"本文围绕 `{topic}` 组织内容，采用“概念 -> 原理 -> 误区 -> 练习”的生成链路展开。",
            "",
            "## 生成大纲",
            *[f"- {plan.title}: {plan.objective}" for plan in section_plans],
        ]

        for section_index, (plan, generated_section) in enumerate(
            zip(section_plans, generated_sections.sections, strict=False),
            start=1,
        ):
            lines.extend(
                [
                    "",
                    f"## {generated_section.title or plan.title}",
                    generated_section.body,
                    "",
                    "### 学习提示",
                    *generated_section.tips,
                    "",
                    "### 引用依据",
                    *generated_section.citations,
                ]
            )

        lines.extend(["", "## 参考来源", *self._render_source_catalog(sources)])
        return "\n".join(lines)

    def render_section_paragraph(
        self,
        *,
        plan: SectionPlan | dict[str, Any],
        snapshot: SystemSnapshot | dict[str, Any],
        topic: str,
        section_index: int,
    ) -> str:
        snapshot_student_level = self._snapshot_value(snapshot, "student_level")
        snapshot_gaps = self._snapshot_list(snapshot, "knowledge_gaps")
        paragraph_by_section = {
            1: (
                f"`{topic}` 是当前知识点中最先需要建立的概念锚点。"
                f" 对于 `{snapshot_student_level}` 水平的学生，建议先回答“它是什么、解决什么问题、和相邻概念有什么区别”。"
            ),
            2: (
                f"理解 `{topic}` 时，不要只记结论，更要抓住判断条件与使用边界。"
                " 学习时可以结合课程中的典型例题，把原理和题目条件一一对应起来。"
            ),
            3: (
                f"从当前画像看，学生在 `{', '.join(snapshot_gaps) or '暂无明确薄弱点'}` 上更容易出错。"
                f" 因此本节重点解释 `{topic}` 与易混概念之间的边界，以及常见错因。"
            ),
            4: (
                f"完成 `{topic}` 的学习后，建议立即安排小规模练习，并把错题回流到薄弱点记录中。"
                " 先做基础题验证概念，再做综合题训练迁移。"
            ),
        }
        objective = plan.get("objective", "") if isinstance(plan, dict) else plan.objective
        return paragraph_by_section.get(section_index, objective)

    def render_section_tips(
        self,
        *,
        plan: SectionPlan | dict[str, Any],
        snapshot: SystemSnapshot | dict[str, Any],
        section_index: int,
    ) -> list[str]:
        current_chapter = self._snapshot_value(snapshot, "current_chapter")
        tips_by_section = {
            1: [f"- 先用一句话复述概念，再和 `{current_chapter}` 中相邻知识点做区分。"],
            2: ["- 尝试把判断条件写成 2-3 条清单，减少“只会背不会用”的情况。"],
            3: ["- 做题时先圈出限制条件，再判断是否满足使用前提。"],
            4: ["- 练习后及时记录错因，并把错因映射回当前画像薄弱点。"],
        }
        objective = plan.get("objective", "") if isinstance(plan, dict) else plan.objective
        return tips_by_section.get(section_index, [f"- {objective}"])

    def render_section_citations(self, *, plan: SectionPlan | dict[str, Any]) -> list[str]:
        source_titles = (
            plan.get("sourceTitles", [])
            if isinstance(plan, dict)
            else plan.source_titles
        )
        if not source_titles:
            return ["- [来源] 当前为回退内容，待补真实检索来源。"]
        return [
            f"- [来源{index}] {source_title}"
            for index, source_title in enumerate(source_titles, start=1)
        ]

    def _render_source_catalog(self, sources: list[dict[str, Any]]) -> list[str]:
        if not sources:
            return ["- 暂无稳定来源"]
        lines: list[str] = []
        for index, source in enumerate(sources[:5], start=1):
            title = str(source.get("title", "未知来源"))
            channel = str(source.get("channel", "unknown"))
            score = source.get("score", "")
            evidence = source.get("evidence")
            lines.append(f"- [来源{index}] {title} ({channel}:{score})")
            if evidence:
                lines.append(f"  证据说明: {evidence}")
        return lines

    def _build_preview_text(self, section_plans: list[SectionPlan]) -> str:
        if not section_plans:
            return "已生成课程资源"
        return f"已生成结构化文档，共 {len(section_plans)} 个章节"

    def _snapshot_value(
        self,
        snapshot: SystemSnapshot | dict[str, Any],
        key: str,
    ) -> Any:
        if isinstance(snapshot, dict):
            return snapshot.get(key, "")
        return getattr(snapshot, key)

    def _snapshot_list(
        self,
        snapshot: SystemSnapshot | dict[str, Any],
        key: str,
    ) -> list[str]:
        value = self._snapshot_value(snapshot, key)
        return list(value) if isinstance(value, list) else []

    def _build_generation_snapshot(
        self,
        *,
        params: dict[str, Any],
        snapshot: SystemSnapshot,
    ) -> dict[str, Any]:
        profile = params.get("profile", {}) if isinstance(params.get("profile", {}), dict) else {}
        base = asdict(snapshot)
        base["preferred_resource_types"] = list(profile.get("preferredResourceTypes", []))
        base["learning_goal"] = (
            profile.get("learningGoal")
            or (profile.get("currentGoal", {}) or {}).get("shortTerm")
            or ""
        )
        return base

    def _build_reading(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        title = f"{params.get('query', '学习主题')}延伸阅读"
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", [])
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        generated_reading = self.content_chain.generate_reading_asset(
            title=title,
            topic=str(params.get("rewrittenQuery", params.get("query", "主题"))),
            snapshot=generation_snapshot,
            sources=sources,
            fallback_builder=self,
        )
        content = "\n".join([f"# {generated_reading.title}", "", generated_reading.body])
        file_name = self._scoped_file_name("reading_material", "md", params)
        path = self._write_text(file_name, content)
        return GeneratedAsset(
            assetType="READING",
            title=generated_reading.title,
            summary=generated_reading.summary,
            displayMode="MARKDOWN_CARD",
            fileName=file_name,
            localPath=str(path),
            previewText=generated_reading.title,
            mimeType="text/markdown; charset=UTF-8",
            inlineContent=content,
        )

    def _build_slides(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        title = f"{params.get('query', '学习主题')}PPT大纲"
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", [])
        topic = str(params.get("rewrittenQuery", params.get("query", "主题")))
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)

        # ── 尝试使用 MiMo-V2-Omni 生成 PPTX ──
        pptx_bytes = self._generate_pptx_with_omni(
            title=title, topic=topic, snapshot=snapshot, sources=sources
        )
        if pptx_bytes is not None:
            file_name = self._scoped_file_name("slides", "pptx", params)
            path = self._write_bytes(file_name, pptx_bytes)
            slide_count = self._count_pptx_slides(pptx_bytes)
            return GeneratedAsset(
                assetType="SLIDES",
                title=title,
                summary=f"MiMo-V2-Omni 生成的 PPT 演示文稿 ({slide_count} 页)",
                displayMode="DOWNLOAD_CARD",
                fileName=file_name,
                localPath=str(path),
                previewText=f"PPT 演示文稿 · {slide_count} 页 · {topic}",
                mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

        # ── 回退到 LLM + markdown ──
        generated_slides = self.content_chain.generate_slides_asset(
            title=title,
            topic=topic,
            snapshot=generation_snapshot,
            sources=sources,
            fallback_builder=self,
        )
        slide_lines: list[str] = []
        for index, slide in enumerate(generated_slides.slides, start=1):
            slide_lines.extend(
                [
                    f"{index}. {slide.title}",
                    *[f"   - {bullet}" for bullet in slide.bullets],
                    f"   讲解备注: {slide.speaker_notes}",
                ]
            )
        content = "\n".join(
            [
                f"# {generated_slides.title}",
                "",
                *slide_lines,
            ]
        )
        file_name = self._scoped_file_name("slides_outline", "md", params)
        path = self._write_text(file_name, content)
        return GeneratedAsset(
            assetType="SLIDES",
            title=generated_slides.title,
            summary=generated_slides.summary,
            displayMode="SLIDE_CARD",
            fileName=file_name,
            localPath=str(path),
            previewText=generated_slides.title,
        )

    def _generate_pptx_with_omni(
        self,
        *,
        title: str,
        topic: str,
        snapshot: SystemSnapshot,
        sources: list[dict[str, Any]],
    ) -> bytes | None:
        """尝试通过 MiMo-V2-Omni 生成 PPTX；失败时返回 None。"""
        settings = get_settings()
        if not settings.mimo_api_key:
            return None

        try:
            from src.ai_modules.llms.mimo_client import MiMoClient

            source_texts = "\n".join(
                f"- {s.get('title', 'unknown')}: {s.get('evidence', '')[:200]}"
                for s in sources[:4]
            )
            prompt = (
                f"请为教学主题「{topic}」生成一份完整的 PPT 内容，用于 {snapshot.current_course} 课程。\n"
                f"学生水平: {snapshot.student_level}，学习风格: {snapshot.preferred_style}。\n"
                f"参考来源:\n{source_texts}\n\n"
                "请以JSON格式输出，包含以下字段：\n"
                '{{"slides":[{{"slideTitle":"标题","bullets":["要点1","要点2"],"speakerNotes":"讲解备注"}}]}}\n'
                "要求：6-10页幻灯片，每页3-5个要点，speakerNotes用中文写50-100字的讲解说明。"
                "首尾页分别为标题页和总结页。仅输出JSON。"
            )

            client = MiMoClient()
            response = client.omni_chat_sync(
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=4096,
            )
            slide_data = client.extract_json(response)
            if not slide_data:
                return None

            slides = slide_data.get("slides", [])
            if not slides:
                return None

            return self._build_pptx_bytes(
                title=title,
                topic=topic,
                slides=slides,
                course=str(snapshot.current_course),
            )
        except Exception:
            return None

    @staticmethod
    def _build_pptx_bytes(
        *,
        title: str,
        topic: str,
        slides: list[dict[str, Any]],
        course: str,
    ) -> bytes:
        """使用 python-pptx 在内存中构建 PPTX 文件。"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return None

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ── 标题页 ──
        title_slide_layout = prs.slide_layouts[0]  # 标题页布局
        slide = prs.slides.add_slide(title_slide_layout)
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        title_placeholder.text = title
        subtitle_placeholder.text = f"{course}\n{topic}"

        # ── 内容页 ──
        for slide_info in slides:
            slide_title = slide_info.get("slideTitle", slide_info.get("title", ""))
            bullets = slide_info.get("bullets", [])
            speaker_notes_text = slide_info.get("speakerNotes", slide_info.get("speaker_notes", ""))

            bullet_layout = prs.slide_layouts[1]  # 标题 + 内容
            slide = prs.slides.add_slide(bullet_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_title

            # 添加要点
            body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if body_shape and bullets:
                tf = body_shape.text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.level = 0
                    p.font.size = Pt(24)

            # 讲解备注
            if speaker_notes_text:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = str(speaker_notes_text)

        # ── 总结页 ──
        summary_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(summary_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "总结与回顾"
        body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            summary_points = [
                f"主题: {topic}",
                f"课程: {course}",
                f"共 {len(slides)} 个内容页",
                "请结合课堂讨论加深理解",
            ]
            for i, point in enumerate(summary_points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = point
                p.font.size = Pt(24)

        import io
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    @staticmethod
    def _count_pptx_slides(pptx_bytes: bytes) -> int:
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(pptx_bytes))
            return len(prs.slides)
        except Exception:
            return 0

    def _write_bytes(self, file_name: str, data: bytes) -> Path:
        self.sandbox_root.mkdir(parents=True, exist_ok=True)
        path = self.sandbox_root / file_name
        path.write_bytes(data)
        return path

    def _build_mindmap(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        title = f"{params.get('query', '学习主题')}思维导图"
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", [])
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        generated_mindmap = self.content_chain.generate_mindmap_asset(
            title=title,
            topic=str(params.get("rewrittenQuery", params.get("query", "主题"))),
            snapshot=generation_snapshot,
            sources=sources,
            fallback_builder=self,
        )
        mermaid = self._render_mindmap_mermaid(generated_mindmap)
        file_name = self._scoped_file_name("mindmap", "mmd", params)
        path = self._write_text(file_name, mermaid)
        return GeneratedAsset(
            assetType="MINDMAP",
            title=generated_mindmap.title,
            summary=generated_mindmap.summary,
            displayMode="INLINE_MERMAID",
            fileName=file_name,
            localPath=str(path),
            previewText=generated_mindmap.title,
            mimeType="text/plain; charset=UTF-8",
            inlineContent=mermaid,
        )

    def _render_mindmap_mermaid(self, generated_mindmap: GeneratedMindMap) -> str:
        lines = ["mindmap", f'  root["{self._escape_mermaid_label(generated_mindmap.root)}"]']
        node_index = 0

        def append_nodes(nodes: list[Any], depth: int) -> None:
            nonlocal node_index
            indent = "  " * depth
            for node in nodes:
                node_index += 1
                lines.append(
                    f'{indent}node_{node_index}["{self._escape_mermaid_label(str(node.name))}"]'
                )
                append_nodes(node.children, depth + 1)

        append_nodes(generated_mindmap.children, 2)
        return "\n".join(lines) + "\n"

    @staticmethod
    def _escape_mermaid_label(label: str) -> str:
        return (
            str(label)
            .replace("\\", "\\\\")
            .replace('"', '\\"')
            .replace("\r", " ")
            .replace("\n", " ")
            .strip()
        )

    def _build_code(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        title = f"{params.get('query', '学习主题')}代码案例"
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", [])
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        generated_code = self.content_chain.generate_code_asset(
            title=title,
            topic=str(params.get("rewrittenQuery", params.get("query", "主题"))),
            snapshot=generation_snapshot,
            sources=sources,
            fallback_builder=self,
        )
        code_suffix = self._code_file_suffix(generated_code.language)
        file_name = self._scoped_file_name("code_case", code_suffix, params)
        path = self._write_text(file_name, generated_code.code)
        return GeneratedAsset(
            assetType="CODE",
            title=generated_code.title,
            summary=generated_code.summary,
            displayMode="INLINE_CODE",
            fileName=file_name,
            localPath=str(path),
            previewText=generated_code.title,
            mimeType="text/plain; charset=UTF-8",
            inlineContent=generated_code.code,
            language=generated_code.language,
            explanation=generated_code.explanation,
        )

    def _build_video(self, *, params: dict, snapshot: SystemSnapshot) -> GeneratedAsset:
        topic = str(params.get("topic") or params.get("query") or "教学主题")
        task_id = str(params.get("taskId") or "video-task")
        style = str(params.get("style") or "hybrid")
        duration_target = self._normalize_duration_seconds(params.get("duration") or params.get("durationTarget"))
        retrieval = params.get("retrievalResult", {})
        sources = retrieval.get("documents", []) if isinstance(retrieval, dict) else []
        generation_snapshot = self._build_generation_snapshot(params=params, snapshot=snapshot)
        script_payload = self.content_chain.generate_video_script(
            title=f"{topic}教学视频",
            topic=topic,
            snapshot=generation_snapshot,
            sources=sources,
            duration_seconds=duration_target,
            style=style,
            fallback_builder=self,
        )
        task_dir = self.sandbox_root / f"video_{self._safe_task_id(task_id)}"
        task_dir.mkdir(parents=True, exist_ok=True)

        script_json_path = task_dir / "script.json"
        script_text_path = task_dir / "script.txt"
        audio_path = task_dir / "speech.mp3"
        thumbnail_path = task_dir / "thumbnail.svg"

        script_json_path.write_text(
            json.dumps(script_payload.model_dump(by_alias=True), ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        script_text_path.write_text(script_payload.full_text, encoding="utf-8")

        # TTS 必须朗读最终脚本，而不是原始检索证据。
        tts_audio_bytes = params.get("tts_audio_bytes")
        if not tts_audio_bytes or not isinstance(tts_audio_bytes, bytes) or len(tts_audio_bytes) < 100:
            from src.ai_modules.llms.mimo_client import MiMoClient

            try:
                mimo_client = MiMoClient()
                tts_audio_bytes = mimo_client.synthesize_speech_sync(
                    text=script_payload.full_text[:1600],
                    style_description="用清晰自然的语速播报，声音沉稳专业，适合教学场景",
                    voice="mimo_default",
                    audio_format="mp3",
                )
            except Exception as exc:
                raise RuntimeError(
                    f"Video TTS generation failed: {type(exc).__name__}: {exc}"
                ) from exc
            params["tts_audio_bytes"] = tts_audio_bytes
        audio_path.write_bytes(tts_audio_bytes)

        # 缩略图
        thumbnail_path.write_text(
            self._build_video_thumbnail_svg(
                title=script_payload.title,
                topic=topic,
                style=style,
            ),
            encoding="utf-8",
        )

        artifact = VideoSandboxArtifact(
            taskDir=str(task_dir),
            scriptJsonPath=str(script_json_path),
            scriptTextPath=str(script_text_path),
            audioPath=str(audio_path),
            finalVideoPath=None,
            thumbnailPath=str(thumbnail_path),
            durationSeconds=script_payload.total_duration,
            videoStyle=style,
            previewText=script_payload.full_text[:100],
            summaryText=f"{topic} 教学视频脚本与语音已生成，等待浏览器本地渲染。",
            audioBase64=base64.b64encode(tts_audio_bytes).decode("utf-8"),
            audioFormat="mp3",
            avatarDataUrl="/dh_live/assets/combined_data.json.gz",
        )
        settings = get_settings()
        params["videoSandboxArtifact"] = artifact.model_dump(by_alias=True)
        params["videoGenerationTask"] = VideoGenerationTaskPayload(
            status="completed",
            title=script_payload.title,
            topic=topic,
            script=script_payload,
            durationSeconds=artifact.duration_seconds,
            videoStyle=style,
            ttsProvider=settings.tts_provider,
            avatarProvider="browser_dh_live_mini",
            generationParams={
                "durationTarget": script_payload.total_duration,
                "style": style,
            },
        ).model_dump(by_alias=True)
        return GeneratedAsset(
            assetType="VIDEO",
            title=script_payload.title,
            summary=artifact.summary_text,
            displayMode="VIDEO_PLAYER",
            fileName="browser-rendered.webm",
            localPath=None,
            previewText=artifact.preview_text,
            mimeType="video/webm",
            thumbnailPath=str(thumbnail_path),
            thumbnailFileName="thumbnail.svg",
            thumbnailMimeType="image/svg+xml",
            durationSeconds=artifact.duration_seconds,
            videoStyle=style,
            knowledgePoint=topic,
        )

    def _safe_task_id(self, task_id: str) -> str:
        return "".join(ch if ch.isalnum() or ch in ("-", "_") else "_" for ch in task_id)

    def _normalize_duration_seconds(self, value: Any) -> int:
        try:
            duration = int(value)
        except (TypeError, ValueError):
            duration = 60
        return max(15, min(180, duration))

    def _build_video_thumbnail_svg(self, *, title: str, topic: str, style: str) -> str:
        style_label = {
            "talking_head": "数字人讲解",
            "animation": "动画演示",
            "hybrid": "混合讲解",
        }.get(style, "教学视频")
        safe_title = self._escape_svg_text(title)
        safe_topic = self._escape_svg_text(topic)
        safe_style = self._escape_svg_text(style_label)
        return (
            '<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">'
            '<defs><linearGradient id="bg" x1="0" y1="0" x2="1" y2="1">'
            '<stop offset="0%" stop-color="#0f172a"/><stop offset="100%" stop-color="#2563eb"/>'
            '</linearGradient></defs>'
            '<rect width="1280" height="720" fill="url(#bg)"/>'
            '<rect x="64" y="64" width="1152" height="592" rx="32" fill="rgba(15,23,42,0.28)" stroke="#93c5fd" stroke-width="2"/>'
            '<text x="100" y="160" fill="#cbd5e1" font-size="30" font-family="Segoe UI, Arial, sans-serif">AI Teaching Video</text>'
            f'<text x="100" y="268" fill="#ffffff" font-size="58" font-family="Segoe UI, Arial, sans-serif">{safe_title}</text>'
            f'<text x="100" y="352" fill="#dbeafe" font-size="34" font-family="Segoe UI, Arial, sans-serif">知识点: {safe_topic}</text>'
            f'<text x="100" y="412" fill="#bfdbfe" font-size="30" font-family="Segoe UI, Arial, sans-serif">风格: {safe_style}</text>'
            '<circle cx="1070" cy="360" r="84" fill="#ffffff" fill-opacity="0.9"/>'
            '<polygon points="1042,314 1042,406 1118,360" fill="#2563eb"/>'
            "</svg>"
        )

    def _escape_svg_text(self, value: str) -> str:
        return (
            value.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
        )

    def build_fallback_reading_asset(
        self,
        *,
        title: str,
        topic: str,
        snapshot: dict[str, Any],
        sources: list[dict[str, Any]],
    ) -> GeneratedTextAsset:
        del title, topic, snapshot, sources
        raise RuntimeError("Local fallback reading asset generation is disabled; use LLM generation only.")

    def build_fallback_slides_asset(
        self,
        *,
        title: str,
        topic: str,
        snapshot: dict[str, Any],
        sources: list[dict[str, Any]],
    ) -> GeneratedSlideDeck:
        del title, topic, snapshot, sources
        raise RuntimeError("Local fallback slides generation is disabled; use LLM generation only.")

    def build_fallback_mindmap_asset(
        self,
        *,
        title: str,
        topic: str,
        snapshot: dict[str, Any],
        sources: list[dict[str, Any]],
    ) -> GeneratedMindMap:
        del title, topic, snapshot, sources
        raise RuntimeError("Local fallback mindmap generation is disabled; use LLM generation only.")

    def build_fallback_code_asset(
        self,
        *,
        title: str,
        topic: str,
        snapshot: dict[str, Any],
        sources: list[dict[str, Any]],
    ) -> GeneratedCodeAsset:
        del title, topic, snapshot, sources
        raise RuntimeError("Local fallback code generation is disabled; use LLM generation only.")

    def _write_text(self, file_name: str, content: str) -> Path:
        self.sandbox_root.mkdir(parents=True, exist_ok=True)
        path = self.sandbox_root / file_name
        path.write_text(content, encoding="utf-8")
        return path

    @staticmethod
    def _code_file_suffix(language: str | None) -> str:
        normalized = (language or "").strip().lower()
        suffix_by_language = {
            "python": "py",
            "py": "py",
            "javascript": "js",
            "js": "js",
            "typescript": "ts",
            "ts": "ts",
            "java": "java",
            "rust": "rs",
            "rs": "rs",
            "go": "go",
            "c": "c",
            "cpp": "cpp",
            "c++": "cpp",
            "sql": "sql",
            "html": "html",
            "css": "css",
            "shell": "sh",
            "bash": "sh",
        }
        return suffix_by_language.get(normalized, "txt")

    def _scoped_file_name(self, prefix: str, suffix: str, params: dict[str, Any]) -> str:
        task_id = str(params.get("taskId") or "shared")
        safe_task_id = "".join(ch if ch.isalnum() or ch in ("-", "_") else "_" for ch in task_id)
        return f"{prefix}_{safe_task_id}.{suffix}"
